"""
DocuMint Document Extraction
Extracts text from PDF, Word, PowerPoint, and Excel files.
Uses pypdf for PDF extraction (Python 3.14 compatible)
"""

from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook
from pathlib import Path
from typing import Optional


def extract_text(file_path: str) -> str:
    """
    Extract text from a document based on its file type.
    
    Args:
        file_path: Path to the document file
    
    Returns:
        Extracted text as a string
    
    Raises:
        ValueError: If file type is not supported
    """
    path = Path(file_path)
    suffix = path.suffix.lower()
    
    extractors = {
        '.pdf': extract_pdf,
        '.docx': extract_docx,
        '.doc': extract_docx,
        '.pptx': extract_pptx,
        '.xlsx': extract_xlsx,
        '.xls': extract_xlsx,
    }
    
    extractor = extractors.get(suffix)
    if not extractor:
        raise ValueError(f"Unsupported file type: {suffix}")
    
    return extractor(file_path)


def extract_pdf(file_path: str) -> str:
    """
    Extract text from PDF using pypdf.
    Pure Python library - works with any Python version.
    """
    text_parts = []
    
    try:
        reader = PdfReader(file_path)
        
        for page_num, page in enumerate(reader.pages, start=1):
            page_text = page.extract_text()
            
            if page_text and page_text.strip():
                text_parts.append(f"[Page {page_num}]\n{page_text}")
        
    except Exception as e:
        raise RuntimeError(f"Failed to extract PDF: {str(e)}")
    
    return "\n\n".join(text_parts)


def extract_docx(file_path: str) -> str:
    """Extract text from Word document using python-docx."""
    text_parts = []
    
    try:
        doc = DocxDocument(file_path)
        
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_parts.append(" | ".join(row_text))
        
    except Exception as e:
        raise RuntimeError(f"Failed to extract Word document: {str(e)}")
    
    return "\n\n".join(text_parts)


def extract_pptx(file_path: str) -> str:
    """Extract text from PowerPoint using python-pptx."""
    text_parts = []
    
    try:
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            slide_text.append(" | ".join(row_text))
            
            if slide_text:
                text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
        
    except Exception as e:
        raise RuntimeError(f"Failed to extract PowerPoint: {str(e)}")
    
    return "\n\n".join(text_parts)


def extract_xlsx(file_path: str) -> str:
    """Extract text from Excel using openpyxl."""
    text_parts = []
    
    try:
        wb = load_workbook(file_path, data_only=True)
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = [f"[Sheet: {sheet_name}]"]
            
            for row in sheet.iter_rows():
                row_values = []
                for cell in row:
                    if cell.value is not None:
                        row_values.append(str(cell.value))
                
                if row_values:
                    sheet_text.append(" | ".join(row_values))
            
            if len(sheet_text) > 1:
                text_parts.append("\n".join(sheet_text))
        
        wb.close()
        
    except Exception as e:
        raise RuntimeError(f"Failed to extract Excel: {str(e)}")
    
    return "\n\n".join(text_parts)


def get_file_type(filename: str) -> Optional[str]:
    """Get the file type from filename."""
    suffix = Path(filename).suffix.lower()
    
    type_map = {
        '.pdf': 'pdf',
        '.docx': 'docx',
        '.doc': 'docx',
        '.pptx': 'pptx',
        '.xlsx': 'xlsx',
        '.xls': 'xlsx',
    }
    
    return type_map.get(suffix)


def is_supported_file(filename: str) -> bool:
    """Check if the file type is supported."""
    return get_file_type(filename) is not None